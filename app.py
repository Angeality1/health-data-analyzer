#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
تطبيق ويب لتحليل بيانات الوحدات الصحية
Health Data Analyzer - Web Application v3.0
مبني على النسخة Desktop المحسّنة
"""

import streamlit as st
import pandas as pd
import io
from datetime import datetime
from pathlib import Path
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.util import Pt

# Page configuration
st.set_page_config(
    page_title="برنامج تحليل بيانات الوحدات الصحية",
    page_icon="🏥",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Custom CSS for RTL and styling
st.markdown("""
<style>
    .main {
        direction: rtl;
        text-align: right;
    }
    h1, h2, h3 {
        text-align: center;
        color: #1f77b4;
    }
    .stButton>button {
        width: 100%;
        background-color: #1f77b4;
        color: white;
        font-size: 18px;
        font-weight: bold;
        padding: 10px;
        border-radius: 5px;
    }
    .success-box {
        padding: 20px;
        background-color: #d4edda;
        border: 1px solid #c3e6cb;
        border-radius: 5px;
        margin: 10px 0;
        text-align: center;
    }
</style>
""", unsafe_allow_html=True)

# Initialize session state
if 'data_loaded' not in st.session_state:
    st.session_state.data_loaded = False
if 'combined_data' not in st.session_state:
    st.session_state.combined_data = None
if 'health_units' not in st.session_state:
    st.session_state.health_units = []

def merge_related_units(df):
    """Merge related health units"""
    unit_mapping = {
        'وحدة بلهارسيا كفر حجازى': 'وحدة كفر حجازى',
        'وحدة بلهارسيا الابشيط': 'وحدة طب أسرة  الابشيط',
        'وحدة طب أسرة  شويقى سعد': 'وحدة عزبة شويقى سعد',
        'مستشفى تكامل الدواخلية': 'مركز طب اسرة الدواخلية',
        'مستشفى تكامل المعتمدية': 'مركز طب اسرة المعتمدية',
        'مستشفى تكامل صفط تراب': 'مركز طب اسرة صفط تراب',
        'مستشفى تكامل العامرية': 'مركز طب اسرة العامرية',
        'مستشفى تكامل الهياتم': 'مركز طب اسرة الهياتم',
        'وحدة بلهارسيا القيراطية': 'وحدة القيراطية'
    }
    df['HealthOffice'] = df['HealthOffice'].replace(unit_mapping)
    return df

def load_excel_files(uploaded_files):
    """Load and process Excel files"""
    required_cols = ['ResidenceDistrict', 'HealthOffice', 'Sex', 'disease', 
                     'status', 'Outcome', 'Department', 'AgeGroup', 'Month', 'Week']
    
    dataframes = []
    for file in uploaded_files:
        try:
            xl = pd.ExcelFile(file)
            
            # Find correct sheet name (handle trailing spaces)
            sheet_name = None
            for sname in xl.sheet_names:
                if sname.strip() == 'محل سكن':
                    sheet_name = sname
                    break
            
            if not sheet_name:
                st.error(f"❌ الملف {file.name} لا يحتوي على شيت 'محل سكن'")
                return None
            
            df = pd.read_excel(xl, sheet_name=sheet_name)
            
            # Verify columns
            missing = [c for c in required_cols if c not in df.columns]
            if missing:
                st.error(f"❌ الملف {file.name} ينقصه الأعمدة: {', '.join(missing)}")
                return None
            
            dataframes.append(df)
            
        except Exception as e:
            st.error(f"❌ خطأ في قراءة الملف {file.name}: {str(e)}")
            return None
    
    # Combine and process
    combined = pd.concat(dataframes, ignore_index=True)
    combined = merge_related_units(combined)
    district_data = combined[combined['ResidenceDistrict'] == 'محلة ثانى'].copy()
    
    return district_data

def fill_cover_slide(prs, unit_name, data):
    """Fill cover slide with improved formatting"""
    slide = prs.slides[0]
    
    min_date = data['reptdate_gen'].min() if 'reptdate_gen' in data.columns else None
    max_date = data['reptdate_gen'].max() if 'reptdate_gen' in data.columns else None
    
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text
            
            if 'وحده الوبائيات و الترصد' in text:
                shape.text_frame.clear()
                p = shape.text_frame.paragraphs[0]
                p.text = f"الإدارة الصحية بالمحلة الكبرى ثان\nالقسم الوقائي\nوحده الوبائيات و الترصد بوحده\n{unit_name}"
                p.font.name = 'Times New Roman (Headings)'
                p.font.size = Pt(40)
                p.font.bold = True
                p.alignment = 1
                
            elif 'الادارة الصحية ثان' in text or 'من 1/1' in text:
                if pd.notna(min_date) and pd.notna(max_date):
                    start = pd.to_datetime(min_date).strftime("%d/%m/%Y")
                    end = pd.to_datetime(max_date).strftime("%d/%m/%Y")
                    
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0]
                    p.text = f"الإدارة الصحية بالمحلة الكبرى ثان\nمن {start} إلى {end}"
                    p.font.name = 'Times New Roman (Headings)'
                    p.font.size = Pt(32)
                    p.font.bold = True
                    p.alignment = 1

def fill_info_slide(prs, unit_name, data, monitor, supervisor, director):
    """Fill info slide with improved formatting"""
    slide = prs.slides[1]
    
    min_date = data['reptdate_gen'].min() if 'reptdate_gen' in data.columns else None
    max_date = data['reptdate_gen'].max() if 'reptdate_gen' in data.columns else None
    
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text
            
            if 'وحده الوبائيات و الترصد بالوحده الصحية' in text:
                shape.text_frame.clear()
                p = shape.text_frame.paragraphs[0]
                p.text = f"وحده الوبائيات و الترصد\n{unit_name}"
                p.font.name = 'Arial (Body)'
                p.font.size = Pt(18)
                p.font.bold = True
                p.alignment = 1
            
            elif 'الوحدة الصحية ب' in text:
                shape.text_frame.clear()
                p = shape.text_frame.paragraphs[0]
                p.text = unit_name
                p.font.name = 'Arial (Body)'
                p.font.size = Pt(28)
                p.font.bold = True
                p.alignment = 1
            
            elif 'الفترة من' in text:
                if pd.notna(min_date) and pd.notna(max_date):
                    start_date = pd.to_datetime(min_date)
                    end_date = pd.to_datetime(max_date)
                    
                    months_ar = {
                        1: 'يناير', 2: 'فبراير', 3: 'مارس', 4: 'أبريل',
                        5: 'مايو', 6: 'يونيو', 7: 'يوليو', 8: 'أغسطس',
                        9: 'سبتمبر', 10: 'أكتوبر', 11: 'نوفمبر', 12: 'ديسمبر'
                    }
                    
                    start_str = f"{start_date.day} {months_ar[start_date.month]} {start_date.year}"
                    end_str = f"{end_date.day} {months_ar[end_date.month]} {end_date.year}"
                    
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0]
                    p.text = f"الفترة من {start_str} وحتى {end_str}"
                    p.font.name = 'Arial (Body)'
                    p.font.size = Pt(28)
                    p.font.bold = True
                    p.alignment = 1
    
    # Update staff names by shape index
    if len(slide.shapes) > 8 and slide.shapes[8].has_text_frame:
        shape = slide.shapes[8]
        shape.text_frame.clear()
        p = shape.text_frame.paragraphs[0]
        p.text = f"أ./ {monitor}" if monitor else "أ./"
        p.font.name = 'Arial (Body)'
        p.font.size = Pt(18)
        p.font.bold = True
    
    if len(slide.shapes) > 9 and slide.shapes[9].has_text_frame:
        shape = slide.shapes[9]
        shape.text_frame.clear()
        p = shape.text_frame.paragraphs[0]
        p.text = f"د./ {supervisor}" if supervisor else "د./"
        p.font.name = 'Arial (Body)'
        p.font.size = Pt(18)
        p.font.bold = True
    
    if len(slide.shapes) > 10 and slide.shapes[10].has_text_frame:
        shape = slide.shapes[10]
        shape.text_frame.clear()
        p = shape.text_frame.paragraphs[0]
        p.text = f"د./ {director}" if director else "د./"
        p.font.name = 'Arial (Body)'
        p.font.size = Pt(18)
        p.font.bold = True

def fill_demographic_slide(prs, unit_name, population, villages):
    """Fill demographic slide with improved formatting"""
    slide = prs.slides[2]
    
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text
            
            if 'البيانات الديموغرافية بالوحده الصحية' in text:
                shape.text_frame.clear()
                p = shape.text_frame.paragraphs[0]
                p.text = f"البيانات الديموغرافية بـ{unit_name} بالإدارة الصحية المحلة ثان"
                p.font.name = 'Arial (Body)'
                p.font.size = Pt(20)
                p.font.bold = True
                p.alignment = 1
            
            elif 'عــــدد السكـــان التقديري' in text:
                shape.text_frame.clear()
                
                p1 = shape.text_frame.paragraphs[0]
                p1.text = f"عــــدد السكـــان التقديري : {population if population else '______'}"
                p1.font.name = 'Arial (Body)'
                p1.font.size = Pt(18)
                p1.font.bold = True
                
                p2 = shape.text_frame.add_paragraph()
                p2.text = f"\nالقرى والعزب التابعة : {villages if villages else '______'}"
                p2.font.name = 'Arial (Body)'
                p2.font.size = Pt(20)
                p2.font.bold = True

def update_chart(slide, categories, values, series_name):
    """Update chart data"""
    for shape in slide.shapes:
        if shape.has_chart:
            chart = shape.chart
            chart_data = CategoryChartData()
            chart_data.categories = categories
            chart_data.add_series(series_name, values)
            chart.replace_data(chart_data)
            break

def generate_report(data, unit_name, population, villages, monitor, supervisor, director):
    """Generate PowerPoint report with all improvements"""
    
    # Load template
    template_path = Path(__file__).parent / 'template.pptx'
    if not template_path.exists():
        st.error("❌ لم يتم العثور على قالب PowerPoint")
        return None
    
    try:
        prs = Presentation(str(template_path))
        
        # Filter Q1 data (including weeks 52, 53)
        q1_data = data[
            (data['Month'].isin([1, 2, 3])) | 
            (data['Week'].isin([52, 53]))
        ].copy()
        
        if len(q1_data) == 0:
            st.warning("⚠️ لا توجد بيانات للربع الأول")
            return None
        
        # Fill slides with improved formatting
        fill_cover_slide(prs, unit_name, q1_data)
        fill_info_slide(prs, unit_name, q1_data, monitor, supervisor, director)
        fill_demographic_slide(prs, unit_name, population, villages)
        
        # Fill all charts
        # Gender chart (slide 4)
        gender_counts = q1_data['Sex'].value_counts()
        categories = ['ذكر', 'أنثى']
        values = [int(gender_counts.get('ذكر', 0)), int(gender_counts.get('أنثى', 0))]
        update_chart(prs.slides[3], categories, values, 'العدد')
        
        # Department chart (slide 5)
        dept_counts = q1_data['Department'].value_counts()
        update_chart(prs.slides[4], dept_counts.index.tolist()[:10], 
                    [int(v) for v in dept_counts.values.tolist()[:10]], 'العدد')
        
        # Outcome chart (slide 6)
        outcome_counts = q1_data['Outcome'].value_counts()
        update_chart(prs.slides[5], outcome_counts.index.tolist()[:10],
                    [int(v) for v in outcome_counts.values.tolist()[:10]], 'العدد')
        
        # Status chart (slide 7)
        status_counts = q1_data['status'].value_counts()
        update_chart(prs.slides[6], status_counts.index.tolist()[:10],
                    [int(v) for v in status_counts.values.tolist()[:10]], 'العدد')
        
        # Disease chart (slide 8)
        disease_counts = q1_data['disease'].value_counts()
        update_chart(prs.slides[7], disease_counts.index.tolist()[:10],
                    [int(v) for v in disease_counts.values.tolist()[:10]], 'العدد')
        
        # Age chart (slide 9)
        age_counts = q1_data['AgeGroup'].value_counts()
        age_order = ['[0-2]', '[>2-5]', '[>5-15]', '[>15-35]', '[>35-50]', '[>50-65]', '[>65]']
        categories = [ag for ag in age_order if ag in age_counts.index]
        values = [int(age_counts.get(ag, 0)) for ag in categories]
        update_chart(prs.slides[8], categories, values, 'العدد')
        
        # Weekly chart (slide 10)
        week_counts = q1_data['Week'].value_counts().sort_index()
        all_weeks = []
        if 52 in week_counts.index:
            all_weeks.append(52)
        if 53 in week_counts.index:
            all_weeks.append(53)
        all_weeks.extend(range(1, 14))
        categories = [f"أسبوع {w}" for w in all_weeks if w in week_counts.index]
        values = [int(week_counts.get(w, 0)) for w in all_weeks if w in week_counts.index]
        update_chart(prs.slides[9], categories, values, 'العدد')
        
        # Monthly chart (slide 11)
        month_counts = q1_data['Month'].value_counts().sort_index()
        month_names = {1: 'يناير', 2: 'فبراير', 3: 'مارس'}
        categories = [month_names.get(m, f"شهر {m}") for m in [1, 2, 3] if m in month_counts.index]
        values = [int(month_counts.get(m, 0)) for m in [1, 2, 3] if m in month_counts.index]
        update_chart(prs.slides[10], categories, values, 'العدد')
        
        # Save to BytesIO
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        
        return output
        
    except Exception as e:
        st.error(f"❌ خطأ في إنشاء التقرير: {str(e)}")
        return None

# Main App
def main():
    # Header
    st.markdown("<h1>🏥 برنامج تحليل بيانات الوحدات الصحية</h1>", unsafe_allow_html=True)
    st.markdown("<h3>الإدارة الصحية بالمحلة الكبرى ثان - القسم الوقائي</h3>", unsafe_allow_html=True)
    st.markdown("---")
    
    # Sidebar
    with st.sidebar:
        st.markdown("### 📋 معلومات التطبيق")
        st.info("""
        **النسخة:** 3.0 (Web)
        
        **المميزات:**
        - ✅ يعمل على الجوال والكمبيوتر
        - ✅ دمج الوحدات المرتبطة
        - ✅ معالجة الأسابيع 52 و 53
        - ✅ تنسيق محسّن للشرائح
        - ✅ تقارير احترافية
        """)
        
        st.markdown("---")
        st.markdown("### 📞 الدعم")
        st.markdown("للمساعدة: راجع دليل الاستخدام")
    
    # Step 1: Upload Files
    st.markdown("### 1️⃣ رفع ملفات Excel")
    uploaded_files = st.file_uploader(
        "اختر ملفات Excel (شهر أو ربع سنوي)",
        type=['xlsx', 'xls'],
        accept_multiple_files=True,
        help="يمكنك رفع ملف واحد (شهري) أو 3 ملفات (ربع سنوي)"
    )
    
    if uploaded_files:
        st.success(f"✅ تم رفع {len(uploaded_files)} ملف")
        
        if st.button("📊 تحميل وتحليل البيانات", key="load_data"):
            with st.spinner("جاري تحميل البيانات..."):
                data = load_excel_files(uploaded_files)
                
                if data is not None:
                    st.session_state.combined_data = data
                    
                    # Get unique units and add "all" option
                    units = sorted(data['HealthOffice'].dropna().unique())
                    units.insert(0, "المحلة ثان (كاملة)")
                    st.session_state.health_units = units
                    st.session_state.data_loaded = True
                    
                    st.success(f"✅ تم تحميل {len(data)} سجل من {len(units)} وحدة صحية")
                    st.balloons()
    
    # Step 2: Select Unit
    if st.session_state.data_loaded:
        st.markdown("---")
        st.markdown("### 2️⃣ اختيار الوحدة الصحية")
        
        selected_unit = st.selectbox(
            "اختر الوحدة:",
            st.session_state.health_units,
            help="اختر وحدة معينة أو 'المحلة ثان (كاملة)' للتحليل الشامل"
        )
        
        # Step 3: Additional Info
        st.markdown("---")
        st.markdown("### 3️⃣ البيانات الإضافية (اختياري)")
        
        col1, col2 = st.columns(2)
        
        with col1:
            population = st.text_input("عدد السكان التقديري", placeholder="مثال: 15000")
            monitor = st.text_input("مراقب الترصد", placeholder="مثال: أحمد محمد")
            supervisor = st.text_input("مسئول الترصد", placeholder="مثال: د. محمد علي")
        
        with col2:
            villages = st.text_input("القرى والعزب التابعة", placeholder="مثال: الهياتم، دنوشر")
            director = st.text_input("مدير الوحدة", placeholder="مثال: د. فاطمة أحمد")
        
        # Step 4: Generate Report
        st.markdown("---")
        st.markdown("### 4️⃣ إنشاء التقرير")
        
        if st.button("📄 إنشاء تقرير PowerPoint", key="generate", type="primary"):
            with st.spinner("جاري إنشاء التقرير... قد يستغرق دقيقة"):
                
                # Get data for selected unit
                if selected_unit == "المحلة ثان (كاملة)":
                    unit_data = st.session_state.combined_data.copy()
                    report_name = "إدارة المحلة ثان كاملة"
                else:
                    unit_data = st.session_state.combined_data[
                        st.session_state.combined_data['HealthOffice'] == selected_unit
                    ].copy()
                    report_name = selected_unit
                
                # Generate report
                report_file = generate_report(
                    unit_data, 
                    report_name,
                    population, 
                    villages,
                    monitor, 
                    supervisor, 
                    director
                )
                
                if report_file:
                    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                    filename = f"تقرير_{report_name.replace('/', '_')}_{timestamp}.pptx"
                    
                    # Download button
                    st.markdown("<div class='success-box'>", unsafe_allow_html=True)
                    st.markdown("### ✅ تم إنشاء التقرير بنجاح!")
                    st.download_button(
                        label="⬇️ تحميل التقرير PowerPoint",
                        data=report_file,
                        file_name=filename,
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )
                    st.markdown("</div>", unsafe_allow_html=True)
                    
                    # Statistics
                    st.markdown("---")
                    st.markdown("### 📊 إحصائيات التقرير")
                    
                    q1_data = unit_data[
                        (unit_data['Month'].isin([1, 2, 3])) | 
                        (unit_data['Week'].isin([52, 53]))
                    ]
                    
                    col1, col2, col3, col4 = st.columns(4)
                    col1.metric("إجمالي السجلات", len(q1_data))
                    col2.metric("عدد الأمراض", q1_data['disease'].nunique())
                    col3.metric("الأقسام", q1_data['Department'].nunique())
                    col4.metric("الفئات العمرية", q1_data['AgeGroup'].nunique())
                    
                    st.success("🎉 يمكنك الآن تحميل التقرير!")

if __name__ == "__main__":
    main()
